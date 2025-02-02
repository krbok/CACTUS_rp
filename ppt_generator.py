from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import logging
from typing import Dict, Optional
import tempfile
import os

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PresentationGenerator:
    """Class to generate PowerPoint presentations from research paper summaries."""
    
    def __init__(self):
        """Initialize presentation settings."""
        try:
            self.prs = Presentation()
            self.title_slide_layout = self.prs.slide_layouts[0]
            self.content_slide_layout = self.prs.slide_layouts[1]
            
            # Define text styles
            self.title_font_size = Pt(32)
            self.subtitle_font_size = Pt(24)
            self.body_font_size = Pt(18)
            
            # Define colors
            self.colors = {
                'black': RGBColor(0, 0, 0),
                'white': RGBColor(255, 255, 255)
            }
        except Exception as e:
            logger.error(f"Error initializing presentation: {e}")
            raise
    
    def _add_title_slide(self, title: str) -> None:
        """Add title slide with minimal styling."""
        try:
            slide = self.prs.slides.add_slide(self.title_slide_layout)
            
            # Add title
            title_shape = slide.shapes.title
            title_shape.text = title
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = self.title_font_size
            paragraph.font.color.rgb = self.colors['black']
            paragraph.alignment = PP_ALIGN.CENTER
            
            # Add subtitle
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = "Research Paper Summary"
            subtitle_frame = subtitle_shape.text_frame
            paragraph = subtitle_frame.paragraphs[0]
            paragraph.font.size = self.subtitle_font_size
            paragraph.font.color.rgb = self.colors['black']
            paragraph.alignment = PP_ALIGN.CENTER
            
        except Exception as e:
            logger.error(f"Error adding title slide: {e}")
            raise
    
    def _add_section_slide(self, title: str, content: str) -> None:
        """Add content slide with clean design."""
        try:
            slide = self.prs.slides.add_slide(self.content_slide_layout)
            
            # Add section title
            title_shape = slide.shapes.title
            title_shape.text = title
            title_frame = title_shape.text_frame
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = self.subtitle_font_size
            paragraph.font.color.rgb = self.colors['black']
            paragraph.alignment = PP_ALIGN.LEFT
            
            # Add content
            content_shape = slide.placeholders[1]
            content_frame = content_shape.text_frame
            
            # Clear existing paragraphs
            for _ in range(len(content_frame.paragraphs) - 1):
                p = content_frame.paragraphs[-1]
                tr = p._element
                tr.getparent().remove(tr)
            
            # Add new content with length limit
            p = content_frame.paragraphs[0]
            p.text = content[:2000]  # Limit content length
            p.font.size = self.body_font_size
            p.font.color.rgb = self.colors['black']
            p.alignment = PP_ALIGN.LEFT
            
            # Add spacing
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            
        except Exception as e:
            logger.error(f"Error adding section slide: {e}")
            raise
    
    def generate_presentation(self, sections: Dict[str, str], output_path: Optional[str] = None) -> str:
        """Generate PowerPoint presentation from paper sections."""
        try:
            # Extract title or use default
            title = sections.get('Title', 'Research Paper Summary').strip()
            if not title:
                title = "Research Paper Summary"
            
            # Add title slide
            self._add_title_slide(title)
            
            # Add content slides in order
            section_order = [
                'Abstract',
                'Introduction',
                'Methods',
                'Results',
                'Discussion',
                'Conclusion'
            ]
            
            for section in section_order:
                if section in sections and sections[section].strip():
                    self._add_section_slide(section, sections[section])
            
            # Handle output path
            if output_path is None:
                temp_dir = tempfile.mkdtemp()
                output_path = os.path.join(temp_dir, "research_summary.pptx")
            
            # Save presentation
            self.prs.save(output_path)
            logger.info(f"Presentation saved to: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error generating presentation: {e}")
            return None

def create_ppt(sections: Dict[str, str], output_path: Optional[str] = None) -> str:
    """Create PowerPoint presentation from paper sections."""
    try:
        generator = PresentationGenerator()
        return generator.generate_presentation(sections, output_path)
    except Exception as e:
        logger.error(f"Error in create_ppt: {e}")
        return None
